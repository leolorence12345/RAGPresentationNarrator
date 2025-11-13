"""
Document Processing Utilities
Handles loading and preprocessing of presentation documents.
"""

from pathlib import Path
from typing import List, Dict
import PyPDF2
from pptx import Presentation


class DocumentProcessor:
    """
    Process various document formats (PDF, PPTX) for RAG pipeline.
    """
    
    @staticmethod
    def extract_text_from_pdf(pdf_path: Path) -> str:
        """
        Extract text from PDF file.
        
        Args:
            pdf_path: Path to PDF file
            
        Returns:
            Extracted text
        """
        text = ""
        try:
            with open(pdf_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
        except Exception as e:
            print(f"Error extracting text from PDF {pdf_path}: {e}")
        return text
    
    @staticmethod
    def extract_text_from_pptx(pptx_path: Path) -> str:
        """
        Extract text from PowerPoint file.
        
        Args:
            pptx_path: Path to PPTX file
            
        Returns:
            Extracted text
        """
        text = ""
        try:
            prs = Presentation(pptx_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            print(f"Error extracting text from PPTX {pptx_path}: {e}")
        return text
    
    @staticmethod
    def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """
        Split text into chunks with overlap.
        
        Args:
            text: Input text
            chunk_size: Size of each chunk
            overlap: Overlap between chunks
            
        Returns:
            List of text chunks
        """
        chunks = []
        start = 0
        while start < len(text):
            end = start + chunk_size
            chunk = text[start:end]
            chunks.append(chunk)
            start = end - overlap
        return chunks

